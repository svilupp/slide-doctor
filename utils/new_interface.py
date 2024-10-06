import gradio as gr
from pptx import Presentation
from typing import List
import json
from models import ExtractedIssue, DetectedIssue, IssueCategory, IssueLocation
import os
import subprocess
import fitz
import base64

merged_dict_global = {}

def convert_pptx_to_pdf(pptx_path, output_folder):
    # Check if LibreOffice is installed
    libreoffice_path = '/Applications/LibreOffice.app/Contents/MacOS/soffice'
    if not os.path.exists(libreoffice_path):
        raise FileNotFoundError("LibreOffice is not installed in the default location.")

    # Prepare the output folder
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Convert the .pptx to .pdf using LibreOffice
    command = [
        libreoffice_path,
        '--headless',
        '--convert-to', 'pdf',
        '--outdir', output_folder,
        pptx_path
    ]

    try:
        subprocess.run(command, check=True)
        print(f"Converted {pptx_path} to PDF successfully.")
    except subprocess.CalledProcessError as e:
        print(f"An error occurred while converting the file: {e}")

def pdf_to_images(pdf_path, output_folder):
    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open the PDF
    doc = fitz.open(pdf_path)

    # Dictionary to store slide ID and Base64 image strings
    slide_images = {}

    # Iterate through each page
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)  # Load the page
        pix = page.get_pixmap()         # Render page to an image

        # Define output path for the image
        img_path = os.path.join(output_folder, f"page_{page_num + 1}.png")
        pix.save(img_path)  # Save the image as PNG

        # Convert image to Base64
        with open(img_path, "rb") as img_file:
            base64_string = base64.b64encode(img_file.read()).decode('utf-8')
            slide_images[page_num + 1] = f"data:image/png;base64,{base64_string}"

    return slide_images

# Combined function to handle both conversions
def convert_pptx_to_images(pptx_file, output_folder):
    # Convert PPTX to PDF
    convert_pptx_to_pdf(pptx_file, output_folder)
    
    # Get the name of the converted PDF
    pdf_path = os.path.join(output_folder, os.path.splitext(os.path.basename(pptx_file))[0] + '.pdf')
    
    # Convert PDF to images and return slide images as a dictionary
    slide_images = pdf_to_images(pdf_path, output_folder)
    
    # Return JSON object containing slide images
    return json.dumps(slide_images, indent=4)


def generate_mock_detected_issues() -> List[DetectedIssue]:
    mock_issues = [
        DetectedIssue(
            extracted_issue=ExtractedIssue(
                issue_description="The pie chart on slide 2 lacks clear labels, making it difficult for the audience to interpret the data accurately.",
                element_location=IssueLocation.BODY_VISUAL,
                element_identification_contains_text=None,
                element_identification_verbatim="Pie chart without labels on slide 2",
                severity="medium"
            ),
            category=IssueCategory.CHART,
            page_id=2,
            file="presentation.pptx"
        ),
        DetectedIssue(
            extracted_issue=ExtractedIssue(
                issue_description="There's a spelling error in the title of slide 2. 'Anual' should be 'Annual'.",
                element_location=IssueLocation.TITLE,
                element_identification_contains_text="Anual Report",
                element_identification_verbatim=None,
                severity="high"
            ),
            category=IssueCategory.SPELL,
            page_id=2,
            file="presentation.pptx"
        ),
        DetectedIssue(
            extracted_issue=ExtractedIssue(
                issue_description="The bullet points on slide 3 are not aligned consistently, creating a visually disorganized appearance.",
                element_location=IssueLocation.BODY_TEXT,
                element_identification_contains_text=None,
                element_identification_verbatim="Misaligned bullet points on slide 3",
                severity="low"
            ),
            category=IssueCategory.ALIGNMENT,
            page_id=3,
            file="presentation.pptx"
        ),
        DetectedIssue(
            extracted_issue=ExtractedIssue(
                issue_description="The bar graph on slide 3 uses colors that are too similar, making it challenging to distinguish between different data sets.",
                element_location=IssueLocation.BODY_VISUAL,
                element_identification_contains_text=None,
                element_identification_verbatim="Bar graph with similar colors on slide 3",
                severity="medium"
            ),
            category=IssueCategory.CHART,
            page_id=3,
            file="presentation.pptx"
        ),
        DetectedIssue(
            extracted_issue=ExtractedIssue(
                issue_description="There's a typo in the footer of slide 3. 'Confidental' should be 'Confidential'.",
                element_location=IssueLocation.FOOTER,
                element_identification_contains_text="Confidental",
                element_identification_verbatim=None,
                severity="high"
            ),
            category=IssueCategory.SPELL,
            page_id=3,
            file="presentation.pptx"
        )
    ]
    
    return mock_issues


# Placeholder sample image path
sample_image_path = "https://via.placeholder.com/150"

def extract_text_from_pptx(file):
    prs = Presentation(file.name)
    
    slides_content = {}
    
    for slide_index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        # Join all text in the slide into a single string
        slide_content = "\n".join(slide_text)
        # print(f"Slide {slide_index + 1} Content:\n{slide_content}\n{'-'*40}")
        slides_content[slide_index + 1] = slide_content  # Slide index as key (1-based)
    
    return json.dumps(slides_content, indent=4)


def create_slide_html(issues_data, merged_dict):
    slides = {}

    # Group issues by slide number (page_id)
    for issue in issues_data:
        slide_number = issue.page_id
        if slide_number not in slides:
            slides[slide_number] = []
        slides[slide_number].append(issue)

    html_content = ""

    # Create the HTML structure for each slide
    for slide_number, issues in slides.items():
        # Access the image path directly from merged_dict
        img_path = merged_dict.get(str(slide_number), {}).get("img_path", sample_image_path)

        html_content += f"<div style='border:1px solid #ddd; padding: 10px; margin: 10px 0;'><h3>Slide {slide_number}</h3>"
        html_content += f"<img src='{img_path}' alt='Slide {slide_number} Preview' style='width:150px;'/>"
        for i, issue in enumerate(issues, 1):
            issue_description = issue.extracted_issue.issue_description
            severity = issue.extracted_issue.severity.capitalize()
            html_content += f"<p><strong>Issue {i}:</strong> {issue_description} (Severity: {severity})</p>"
        html_content += "</div>"

    return html_content



def process_ppt(context_info, ppt_upload):
    # Extract text from the uploaded PowerPoint file
    slides_content = extract_text_from_pptx(ppt_upload)
    img_paths = convert_pptx_to_images(ppt_upload, '../pics')

    # print("slides_content: ", slides_content)
    # print("img_paths: ", img_paths)

    slides_content = json.loads(slides_content)
    img_paths = json.loads(img_paths)

    merged_dict = {}
    for key in slides_content.keys():
        if key in img_paths:
            merged_dict[key] = {
                "img_path": img_paths[key],
                "text": slides_content[key]
            }

    print("merged_dict: ", len(merged_dict))    
    
    # Generate mock issues data
    issues_data = generate_mock_detected_issues()
    
    # Create the HTML content for slide view
    slide_html = create_slide_html(issues_data, merged_dict)
    
    # Return both the summary and the HTML content
    return "Mock summary generated", slide_html

# Building the Gradio interface
with gr.Blocks() as demo:
    gr.Markdown("# Context Brief")

    context_input = gr.Textbox(label="Context Information", placeholder="Enter context for the presentation")
    ppt_upload = gr.File(label="Upload PPT", file_count="single", file_types=[".ppt", ".pptx"])

    generate_button = gr.Button("Generate Page View")

    gr.Markdown("### Summary")
    summary_output = gr.Textbox(label="Summary", lines=4, interactive=False)

    gr.Markdown("### Page View")
    slide_sections_container = gr.HTML("")  # Use HTML component to dynamically render slide sections

    # Generate Page View on button click
    generate_button.click(
        process_ppt,
        inputs=[context_input, ppt_upload],
        outputs=[summary_output, slide_sections_container]
    )

demo.launch()
