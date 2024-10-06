from pptx import Presentation

def extract_text_from_pptx(file: str, include_title_prefix: bool = False) -> dict[str, str]:
    prs = Presentation(file)
    
    slides_content = {}
    
    TITLE_FONT_MIN = 16
    
    for slide_index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs and shape.text_frame.paragraphs[0].runs[0].font.size is not None and shape.text_frame.paragraphs[0].runs[0].font.size >= TITLE_FONT_MIN:  # Assuming 16pt as large font
                    if include_title_prefix:
                        slide_text.append(f"TITLE: {shape.text}")
                    else:
                        slide_text.append(shape.text)
                else:
                    slide_text.append(shape.text)
        # Join all text in the slide into a single string
        content = "\n".join(slide_text)
        slides_content[str(slide_index)] = content  # Slide index as key (0-based)
    
    return slides_content
