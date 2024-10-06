import os
from typing import List
from pydantic import BaseModel, Field
from langchain_mistralai import ChatMistralAI
from langchain_core.messages import HumanMessage, SystemMessage
from PIL import Image
import base64
import io
from enum import Enum
from pptx import Presentation
import json
import asyncio
from tenacity import retry, stop_after_attempt, wait_fixed, retry_if_exception_type
from transformers import pipeline
import requests
from pptx.util import Pt
from pptx.dml.color import RGBColor
from textblob import TextBlob


# Constants for fact-checking
BRAVE_API_KEY = os.getenv("BRAVE_API_KEY")
BRAVE_SEARCH_ENDPOINT = "https://api.search.brave.com/res/v1/web/search"

# Fact-Checking NLP Model
nli_model = pipeline("text-classification", model="facebook/bart-large-mnli")


# Initialize the Mistral client with instructor
api_key = os.getenv("MISTRAL_API_KEY")
model = "mistral-large-latest"


def pptx_to_json(slide_deck_path):
    prs = Presentation(slide_deck_path)
    
    slides_content = []
    
    placeholder_mapping = {
        0: "title",
        1: "subtitle",
        2: "body",
        3: "center_title",
        4: "center_subtitle",
        5: "footer",
        6: "date",
        7: "slide_text",
        8: "header",
        9: "description"
    }
    
    # Iterate over slides
    for slide_index, slide in enumerate(prs.slides):
        slide_data = {"slide_index": slide_index, "content": []}
        
        # Iterate over all shapes in the slide
        for shape in slide.shapes:
            # Check if shape has text content
            if shape.has_text_frame:
                # Check if the shape is a placeholder
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    
                    # Map placeholder to a human-readable type name
                    placeholder_name = placeholder_mapping.get(placeholder_type, "unknown_placeholder")
                    
                    slide_data['content'].append({
                        "type": placeholder_name,
                        "text": shape.text
                    })
                else:
                    # If not a placeholder but has text, consider it a regular shape with text
                    slide_data['content'].append({
                        "type": "shape",
                        "text": shape.text
                    })
            
            # Handle non-text placeholder types (e.g., slide number, date, footer)
            elif shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                
                # Map placeholder to a human-readable type name
                placeholder_name = placeholder_mapping.get(placeholder_type, "unknown_placeholder")
                
                if placeholder_name in ["footer", "date", "slide_number"]:
                    slide_data['content'].append({
                        "type": placeholder_name,
                        "text": shape.text if shape.has_text_frame else ""
                    })
        
        slides_content.append(slide_data)
    print('slides_content length: ', len(slides_content))
    
    return slides_content


# Define the output schema
class ObjectType(Enum):
    TITLE = "title"
    TEXT = "text"
    FOOTER = "footer"
    BODY = "body"

class Issue(BaseModel):
    object_type: ObjectType = Field(description="The type of object, from the enum: title, text, chart, footer, table")
    issue_category: str = Field(description="The category of the issue")
    issue_label: str = Field(description="The label of the issue")
    issue_reason: str = Field(description="Reason behind flagging this as an issue")

class ExtractionResult(BaseModel):
    issues: List[Issue] = Field(description="List of extracted issues")

@retry(
    stop=stop_after_attempt(3),
    wait=wait_fixed(1),
    retry=retry_if_exception_type((ValueError, json.JSONDecodeError))
)
def extract_issues(slides_text: str) -> ExtractionResult:

    system_prompt = """
    You are an AI specialized in identifying issues in PowerPoint presentation screenshots.
    Your task is to carefully examine the provided text which is extracted from slides of a .pptx file and identify any problems or areas for improvement related to the slide's content and overall effectiveness.
    Focus on detecting issues such as:
    1. Clarity and readability of text
    2. Spelling mistakes and Effectiveness
    3. organization of information
    4. Consistency in formatting
    5. Appropriateness of content for the intended audience
    Provide a comprehensive analysis of the issues present in the text per slide
    """

    user_prompt = """
    Analyze this slide text and identify issues. For each issue, provide:
    1. The object type (title, text, footer, table)
    2. The issue category (e.g., "Clarity", "Content")
    3. A brief label for the issue (e.g., "Text is too small", "Text does not look professional")
    """

    llm  = ChatMistralAI(mistral_api_key=api_key, model=model)
    structured_llm = llm.with_structured_output(ExtractionResult)

    prompt=[
        SystemMessage(content=system_prompt),
        HumanMessage(content=[
            {"type": "text", "text": user_prompt},
        ])
    ]

    response = structured_llm.invoke(prompt)
    return response


def brave_search(query, count=5):
    headers = {
        "Accept": "application/json",
        "X-Subscription-Token": BRAVE_API_KEY
    }
    
    params = {
        "q": query,
        "count": count
    }
    
    response = requests.get(BRAVE_SEARCH_ENDPOINT, headers=headers, params=params)
    if response.status_code == 200:
        return response.json()
    else:
        response.raise_for_status()

def fact_check(statement):
    search_results = brave_search(query=statement)
    
    snippet_url_pairs = []
    for result in search_results.get("web", {}).get("results", []):
        if result.get('description'):
            snippet_url_pairs.append((result['description'], result.get('url', 'No URL')))
        if result.get('extra_snippets'):
            for snippet in result['extra_snippets']:
                snippet_url_pairs.append((snippet, result.get('url', 'No URL')))

    label_dict = {"ENTAILMENT": None, "CONTRADICTION": None, "NEUTRAL": None}

    for snippet, url in snippet_url_pairs:
        result = nli_model(f"{statement} </s> {snippet}")[0]
        label = result['label'].upper()
        
        if label in label_dict:
            if label_dict[label] is None or result['score'] > label_dict[label][1]:
                label_dict[label] = (snippet, result['score'], url)

    return label_dict


def fix_issue_on_slide(prs, slide_index, issue: Issue):
    slide = prs.slides[slide_index]
    
    # Loop through all shapes in the slide to find any text to fix
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Handle spelling issues by running a spell checker on the text in the shape
            elif issue.issue_category == "Spelling":
                print("Spell checker called!")
                corrected_text = spell_check_correction(shape.text) 
                shape.text = corrected_text
            
            # Handle consistency issues by setting a uniform font size and color for all text in the shape
            elif issue.issue_category == "Consistency":
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(24)  # Set a consistent font size
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Set text color to black
            else:
                



def spell_check_correction(text):
    blob = TextBlob(text)
    
    # Correct the text
    corrected_text = blob.correct()
    
    return str(corrected_text)



def main():


    slide_deck_path = "../data/03-dickinson-basic.pptx"
    
    slides_json = pptx_to_json(slide_deck_path)

    output_slide_deck_path = "../data/03-dickinson-basic_fixed.pptx"
    
    prs = Presentation(slide_deck_path)

    for idx, slide_json in enumerate(slides_json):
        print(f'SLIDE {idx}')
        print('-------------------------------------')

        for obj in slide_json['content']:
            if len(obj['text'].split()) >= 10:
                txt = obj['text']
                print(f'\nFact Checking statement: {txt}')
                fact_check_results = fact_check(issue.issue_label)
                print('fact_check_results: ', fact_check_results)

        slide_response = extract_issues(slide_json)

        # Pretty print the extracted issues
        print("\nExtracted Issues:")
        for i, issue in enumerate(slide_response.issues, 1):
            print(f"\nIssue {i}:")
            print(f"  Object Type: {issue.object_type}")
            print(f"  Category: {issue.issue_category}")
            print(f"  Label: {issue.issue_label}")
            print(f"  Reason: {issue.issue_reason}")

            fix_issue_on_slide(prs, idx, issue)

    prs.save(output_slide_deck_path)
    print(f"Updated presentation saved at {output_slide_deck_path}")

if __name__ == "__main__":
    main()

