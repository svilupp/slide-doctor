import os
from typing import List
from pydantic import BaseModel, Field
from langchain_mistralai import ChatMistralAI
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation
from enum import Enum
import requests
from transformers import pipeline

# Initialize the Mistral client
api_key = os.getenv("MISTRAL_API_KEY")
model = "mistral-large-latest"

# Constants for fact-checking
BRAVE_API_KEY = os.getenv("BRAVE_API_KEY")
BRAVE_SEARCH_ENDPOINT = "https://api.search.brave.com/res/v1/web/search"

# Fact-Checking NLP Model
nli_model = pipeline("text-classification", model="facebook/bart-large-mnli")

# Function to convert PPTX to JSON
def pptx_to_json(slide_deck_path):
    prs = Presentation(slide_deck_path)
    slides_content = []
    
    placeholder_mapping = {
        0: "title",
        1: "subtitle",
        2: "body",
        3: "center_title",
        4: "center_subtitle",
        5: "footer",
        6: "date",
        7: "slide_text",
        8: "header",
    }
    
    for slide_index, slide in enumerate(prs.slides):
        slide_data = {"slide_index": slide_index, "content": []}
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    placeholder_name = placeholder_mapping.get(placeholder_type, "unknown_placeholder")
                    slide_data['content'].append({"type": placeholder_name, "text": shape.text})
                else:
                    slide_data['content'].append({"type": "shape", "text": shape.text})
            elif shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                placeholder_name = placeholder_mapping.get(placeholder_type, "unknown_placeholder")
                if placeholder_name in ["footer", "date", "slide_number"]:
                    slide_data['content'].append({"type": placeholder_name, "text": shape.text if shape.has_text_frame else ""})
        
        slides_content.append(slide_data)
    return slides_content

# Output schema
class ObjectType(Enum):
    TITLE = "title"
    TEXT = "text"
    FOOTER = "footer"
    BODY = "body"

class Issue(BaseModel):
    object_type: ObjectType = Field(description="The type of object, from the enum: title, text, chart, footer, table")
    issue_category: str = Field(description="The category of the issue")
    issue_label: str = Field(description="The label of the issue")
    issue_reason: str = Field(description="Reason behind flagging this as an issue")
    # fact_check: dict = Field(default=None, description="Fact-checking results, if applicable")

class ExtractionResult(BaseModel):
    issues: List[Issue] = Field(description="List of extracted issues")

# Mistral issue extraction
def extract_issues(slides_text: str) -> ExtractionResult:
    system_prompt = """
    You are an AI specialized in identifying issues in PowerPoint presentation screenshots.
    Your task is to carefully examine the provided text extracted from slides and identify any problems related to content and clarity.
    """
    
    user_prompt = "Analyze this slide text and identify issues based on content and formatting."

    llm  = ChatMistralAI(mistral_api_key=api_key, model=model)
    structured_llm = llm.with_structured_output(ExtractionResult)

    prompt = [
        SystemMessage(content=system_prompt),
        HumanMessage(content=[{"type": "text", "text": user_prompt}])
    ]

    response = structured_llm.invoke(prompt)
    return response

# Fact-checking with Brave search and NLI model
def brave_search(query, count=5):
    headers = {
        "Accept": "application/json",
        "X-Subscription-Token": BRAVE_API_KEY
    }
    
    params = {
        "q": query,
        "count": count
    }
    
    response = requests.get(BRAVE_SEARCH_ENDPOINT, headers=headers, params=params)
    if response.status_code == 200:
        return response.json()
    else:
        response.raise_for_status()

def fact_check(statement):
    search_results = brave_search(query=statement)
    
    snippet_url_pairs = []
    for result in search_results.get("web", {}).get("results", []):
        if result.get('description'):
            snippet_url_pairs.append((result['description'], result.get('url', 'No URL')))
        if result.get('extra_snippets'):
            for snippet in result['extra_snippets']:
                snippet_url_pairs.append((snippet, result.get('url', 'No URL')))

    label_dict = {"ENTAILMENT": None, "CONTRADICTION": None, "NEUTRAL": None}

    for snippet, url in snippet_url_pairs:
        result = nli_model(f"{statement} </s> {snippet}")[0]
        label = result['label'].upper()
        
        if label in label_dict:
            if label_dict[label] is None or result['score'] > label_dict[label][1]:
                label_dict[label] = (snippet, result['score'], url)

    return label_dict

# Main function
def main():
    slide_deck_path = "../data/01-coastal-presentation.pptx"
    slides_json = pptx_to_json(slide_deck_path)

    for idx, slide_json in enumerate(slides_json):

        print("slide_json: ", slide_json)

        
        for obj in slide_json['content']:
            if obj['type']=="slide_text" and len(obj['text'].split()) > 5:
                print("\nFact Checks")
                fact_check_results = fact_check(issue.issue_label)
                print('fact_check_results: ', fact_check_results)


        slide_response = extract_issues(slide_json)

        print("\nExtracted Issues:")
        for i, issue in enumerate(slide_response.issues, 1):
            print(f"\nIssue {i}:")
            print(f"  Object Type: {issue.object_type}")
            print(f"  Category: {issue.issue_category}")
            print(f"  Label: {issue.issue_label}")
            print(f"  Reason: {issue.issue_reason}")

if __name__ == "__main__":
    main()
