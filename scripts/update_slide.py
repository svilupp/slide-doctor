import os
from typing import List
from pydantic import BaseModel, Field
from langchain_mistralai import ChatMistralAI
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import json
from enum import Enum
from tenacity import retry, stop_after_attempt, wait_fixed, retry_if_exception_type
from spellchecker import SpellChecker
from textblob import TextBlob


# Initialize the Mistral client with instructor
api_key = os.getenv("MISTRAL_API_KEY")
model = "mistral-large-latest"

def pptx_to_json(slide_deck_path):
    prs = Presentation(slide_deck_path)
    
    slides_content = []
    
    placeholder_mapping = {
        0: "title",
        1: "subtitle",
        2: "body",
        3: "center_title",
        4: "center_subtitle",
        5: "footer",
        6: "date",
        7: "slide_number",
        8: "header",
    }
    
    # Iterate over slides
    for slide_index, slide in enumerate(prs.slides):
        slide_data = {"slide_index": slide_index, "content": []}
        
        # Iterate over all shapes in the slide
        for shape in slide.shapes:
            # Check if shape has text content
            if shape.has_text_frame:
                # Check if the shape is a placeholder
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    
                    # Map placeholder to a human-readable type name
                    placeholder_name = placeholder_mapping.get(placeholder_type, "unknown_placeholder")
                    
                    slide_data['content'].append({
                        "type": placeholder_name,
                        "text": shape.text
                    })
                else:
                    # If not a placeholder but has text, consider it a regular shape with text
                    slide_data['content'].append({
                        "type": "shape",
                        "text": shape.text
                    })
            
            # Handle non-text placeholder types (e.g., slide number, date, footer)
            elif shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                
                # Map placeholder to a human-readable type name
                placeholder_name = placeholder_mapping.get(placeholder_type, "unknown_placeholder")
                
                if placeholder_name in ["footer", "date", "slide_number"]:
                    slide_data['content'].append({
                        "type": placeholder_name,
                        "text": shape.text if shape.has_text_frame else ""
                    })
        
        slides_content.append(slide_data)
    
    return slides_content

# Define the output schema
class ObjectType(Enum):
    TITLE = "title"
    TEXT = "text"
    FOOTER = "footer"
    BODY = "body"

class Issue(BaseModel):
    object_type: ObjectType
    issue_category: str
    issue_label: str
    issue_reason: str

class ExtractionResult(BaseModel):
    issues: List[Issue]


@retry(
    stop=stop_after_attempt(3),
    wait=wait_fixed(1),
    retry=retry_if_exception_type((ValueError, json.JSONDecodeError))
)
def extract_issues(slides_text: str) -> ExtractionResult:
    system_prompt = """
    You are an AI specialized in identifying issues in PowerPoint presentation slides.
    Your task is to carefully examine the provided text which is extracted from slides of a .pptx file and identify any problems or areas for improvement related to the slide's content and overall effectiveness.
    Focus on detecting issues such as:
    1. Spelling checking
    2. Clarity and readability of text
    3. Organization of information
    4. Consistency in formatting
    5. Appropriateness of content for the intended audience
    Provide a comprehensive analysis of the issues present in the text per slide
    """

    user_prompt = """
    Analyze this slide text and identify issues. For each issue, provide:
    1. The object type (title, text, footer, table)
    2. The issue category (e.g., "Clarity", "Content")
    3. A brief label for the issue (e.g., "Text is too small", "Text does not look professional")
    """

    llm = ChatMistralAI(mistral_api_key=api_key, model=model)
    structured_llm = llm.with_structured_output(ExtractionResult)

    prompt = [
        SystemMessage(content=system_prompt),
        HumanMessage(content=[
            {"type": "text", "text": user_prompt}
        ])
    ]

    response = structured_llm.invoke(prompt)
    return response

def fix_issue_on_slide(prs, slide_index, issue: Issue):
    slide = prs.slides[slide_index]
    
    # Loop through shapes in the slide to find the corresponding object to fix
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Example Fixes: Adjust these based on the issue details
            if issue.object_type == ObjectType.TITLE and issue.issue_category == "Clarity":
                # Increase font size if text is unclear
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(28)
            
            elif issue.issue_category == "Spelling":
                print("spell checker called!")
                corrected_text = spell_check_correction(shape.text) 
                shape.text = corrected_text
            
            elif issue.issue_category == "Consistency":
                # Ensure all titles use the same font size and color
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Set to black color



def spell_check_correction(text):
    blob = TextBlob(text)
    
    # Correct the text
    corrected_text = blob.correct()
    
    return str(corrected_text)



def main():
    slide_deck_path = "../data/03-dickinson-basic.pptx"
    output_slide_deck_path = "../data/03-dickinson-basic_fixed.pptx"
    
    prs = Presentation(slide_deck_path)
    slides_json = pptx_to_json(slide_deck_path)

    for idx, slide_json in enumerate(slides_json):
        print(f'SLIDE {idx}')
        print('-------------------------------------')
        
        slide_response = extract_issues(slide_json)
        
        print("\nExtracted Issues:")
        for i, issue in enumerate(slide_response.issues, 1):
            print(f"\nIssue {i}:")
            print(f"  Object Type: {issue.object_type}")
            print(f"  Category: {issue.issue_category}")
            print(f"  Label: {issue.issue_label}")
            print(f"  Reason: {issue.issue_reason}")

            # Fix the slide issue
            fix_issue_on_slide(prs, idx, issue)

    # Save the updated presentation
    prs.save(output_slide_deck_path)
    print(f"Updated presentation saved at {output_slide_deck_path}")

if __name__ == "__main__":
    main()
